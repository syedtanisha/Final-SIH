import io
from sqlalchemy.orm import Session
from typing import List, Dict, Any
from datetime import datetime
from fastapi import UploadFile, HTTPException, status
from pypdf import PdfReader
import docx
from pptx import Presentation

from ..models.models import (
    User, Competency, UserCompetency, LearningProgressHistory,
    Document, Quiz, QuizQuestion, QuizAttempt, LearningResource
)
from ..schemas.assessment import (
    BaselineAssessmentOut, BaselineQuestion, BaselineQuestionOption,
    BaselineAssessmentSubmit, BaselineAssessmentResultOut,
    QuizGenerateRequest, QuizOut, QuizQuestionOut, QuizQuestionDetailOut,
    QuizSubmitRequest, QuizAttemptResultOut, QuestionResultDetail,
    ProgressSummaryOut, ProgressEventOut, CompetencyProgressCard,
    FinalInterviewReadiness, FinalInterviewCompetency
)
from ..data.seed_data import BASELINE_QUESTIONS, COMMON_CORE_QUESTIONS, DEPARTMENT_BASELINE_BANK, ROLE_SPECIFIC_QUESTION_BANK
from .ai_service import (
    generate_mcqs_from_document_async, generate_grok_quiz_feedback,
    generate_final_interview_questions, generate_department_designation_baseline_questions_async
)
from .catalog_service import resolve_role_benchmarks

USER_BASELINE_CACHE: Dict[int, List[Dict[str, Any]]] = {}

import json
import logging

logger = logging.getLogger(__name__)

def select_role_specific_questions(
    role_cat: str,
    target_count: int,
    core_competencies: List[str],
    exclude_ids: Set[int]
) -> List[Dict[str, Any]]:
    selected: List[Dict[str, Any]] = []
    seen_ids = set(exclude_ids)

    COMPATIBLE_TIERS = {
        "senior": ["mid", "technical"],
        "mid": ["senior", "technical", "junior"],
        "junior": ["mid", "technical"],
        "technical": ["mid", "senior"]
    }

    primary_pool = ROLE_SPECIFIC_QUESTION_BANK.get(role_cat, [])

    # Step 1: Exact Role Tier + Exact Core Competency
    for q in primary_pool:
        if len(selected) >= target_count:
            break
        if q["id"] not in seen_ids and q.get("competency_code") in core_competencies:
            q_copy = dict(q)
            q_copy["is_common"] = False
            q_copy["category"] = "role_specific"
            selected.append(q_copy)
            seen_ids.add(q["id"])

    # Step 2: Exact Role Tier + Related/Compatible Competency
    if len(selected) < target_count:
        for q in primary_pool:
            if len(selected) >= target_count:
                break
            if q["id"] not in seen_ids:
                q_copy = dict(q)
                q_copy["is_common"] = False
                q_copy["category"] = "role_specific"
                selected.append(q_copy)
                seen_ids.add(q["id"])

    # Step 3: Compatible Role Tier + Exact Core Competency
    if len(selected) < target_count:
        adj_tiers = COMPATIBLE_TIERS.get(role_cat, [])
        for tier in adj_tiers:
            if len(selected) >= target_count:
                break
            for q in ROLE_SPECIFIC_QUESTION_BANK.get(tier, []):
                if len(selected) >= target_count:
                    break
                if q["id"] not in seen_ids and q.get("competency_code") in core_competencies:
                    q_copy = dict(q)
                    q_copy["is_common"] = False
                    q_copy["category"] = "role_specific"
                    selected.append(q_copy)
                    seen_ids.add(q["id"])

    # Step 4: Approved General Core Question matching Core Competencies
    if len(selected) < target_count:
        for q in COMMON_CORE_QUESTIONS:
            if len(selected) >= target_count:
                break
            if q["id"] not in seen_ids and q.get("competency_code") in core_competencies:
                q_copy = dict(q)
                q_copy["is_common"] = False
                q_copy["category"] = "role_specific"
                selected.append(q_copy)
                seen_ids.add(q["id"])

    # Step 5: Insufficient coverage logging
    if len(selected) < target_count:
        logger.warning(
            f"[QuestionBankCoverage] Insufficient relevant questions for role tier '{role_cat}'. "
            f"Requested {target_count}, available {len(selected)} relevant questions. "
            f"Preserving assessment competency relevance without adding irrelevant substitution."
        )

    return selected

async def get_baseline_assessment_data(
    user: User = None,
    db: Session = None,
    blueprint: Dict[str, Any] = None
) -> BaselineAssessmentOut:
    questions_out: List[BaselineQuestion] = []
    div_title = "MoSPI Statistical System"
    desig_title = "Statistical Officer"
    div_code = "GENERAL"
    role_cat = "mid"
    resolution_method = "exact"
    core_comps: List[str] = ["STAT_SURVEY", "STAT_COMPUTE", "STAT_NAT_ACC", "STAT_PRICE_IND", "STAT_DATA_GOV"]

    if user:
        role_meta = resolve_role_benchmarks(user.department, user.designation)
        div_title = user.department or "MoSPI"
        desig_title = user.designation or "Statistical Officer"
        div_code = role_meta.get("division_code", "GENERAL")
        role_cat = role_meta.get("role_category", "mid")
        resolution_method = role_meta.get("resolution_method", "exact")
        core_comps = role_meta.get("core_competencies", core_comps)

    bp_common_count = 5
    bp_role_count = 4
    if blueprint:
        bp_common_count = blueprint.get("common_core_count", 5)
        bp_role_count = blueprint.get("role_specific_count", 4)

    raw_qs = []
    existing_assignment = None
    if user and db:
        from ..models.models import BaselineAssignment
        existing_assignment = db.query(BaselineAssignment).filter(
            BaselineAssignment.user_id == user.id,
            BaselineAssignment.status == "assigned"
        ).order_by(BaselineAssignment.assigned_at.desc()).first()

    all_known_questions = {}
    for q in COMMON_CORE_QUESTIONS:
        all_known_questions[q["id"]] = dict(q)
    for q_list in ROLE_SPECIFIC_QUESTION_BANK.values():
        for q in q_list:
            all_known_questions[q["id"]] = dict(q)
    for q_list in DEPARTMENT_BASELINE_BANK.values():
        for q in q_list:
            all_known_questions[q["id"]] = dict(q)

    if existing_assignment:
        try:
            assigned_ids = json.loads(existing_assignment.assigned_question_ids)
            for qid in assigned_ids:
                if qid in all_known_questions:
                    raw_qs.append(all_known_questions[qid])
        except Exception as e:
            logger.warning(f"[BaselinePersistence] Error parsing assigned_question_ids: {e}")

    if not raw_qs:
        # Generate 100% Designation and Department-tailored baseline assessment
        target_total_questions = 9
        if blueprint and blueprint.get("total_questions"):
            target_total_questions = blueprint.get("total_questions")

        selected_qs = []
        selected_ids = set()

        # 1. Department-Specific Baseline Questions (e.g. NAD, FOD, PSD, SDRD, DQDD, DES, POLICY)
        dept_pool = DEPARTMENT_BASELINE_BANK.get(div_code, [])
        for q in dept_pool:
            if len(selected_qs) >= 5:  # Up to 5 questions directly from department division bank
                break
            if q["id"] not in selected_ids:
                q_copy = dict(q)
                q_copy["is_common"] = False
                q_copy["category"] = f"department_{div_code.lower()}"
                selected_qs.append(q_copy)
                selected_ids.add(q["id"])

        # 2. Designation Role-Tier Specific Baseline Questions (e.g. senior, mid, junior, technical)
        role_pool = ROLE_SPECIFIC_QUESTION_BANK.get(role_cat, [])
        for q in role_pool:
            if len(selected_qs) >= target_total_questions:
                break
            if q["id"] not in selected_ids:
                q_copy = dict(q)
                q_copy["is_common"] = False
                q_copy["category"] = f"designation_{role_cat}"
                selected_qs.append(q_copy)
                selected_ids.add(q["id"])

        # 3. Fallback to Core Competency Matching if more questions needed
        if len(selected_qs) < target_total_questions:
            needed = target_total_questions - len(selected_qs)
            additional_qs = select_role_specific_questions(
                role_cat=role_cat,
                target_count=needed,
                core_competencies=core_comps,
                exclude_ids=selected_ids
            )
            for q in additional_qs:
                q["is_common"] = False
                selected_qs.append(q)
                selected_ids.add(q["id"])

        raw_qs = selected_qs[:target_total_questions]

        # Persist Assignment to Database
        if user and db:
            from ..models.models import BaselineAssignment
            new_assignment = BaselineAssignment(
                assignment_id=f"mospi-assign-{user.id}-{int(datetime.utcnow().timestamp())}",
                user_id=user.id,
                designation=desig_title,
                department=div_title,
                role_tier=role_cat,
                resolution_method=resolution_method,
                blueprint_version="v2_tailored",
                assigned_question_ids=json.dumps([q["id"] for q in raw_qs]),
                total_questions=len(raw_qs),
                status="assigned"
            )
            db.add(new_assignment)
            db.commit()
            db.refresh(new_assignment)
            existing_assignment = new_assignment

    if user:
        USER_BASELINE_CACHE[user.id] = raw_qs

    assessed_comps = list(set(q["competency_code"] for q in raw_qs))

    for q in raw_qs:
        options = [BaselineQuestionOption(key=opt["key"], text=opt["text"]) for opt in q["options"]]
        questions_out.append(
            BaselineQuestion(
                id=q["id"],
                competency_code=q["competency_code"],
                competency_name=q["competency_name"],
                domain=q["domain"],
                question_text=q["question_text"],
                options=options,
                difficulty=q.get("difficulty", "Intermediate"),
                is_common=q.get("is_common", False),
                category=q.get("category", "common"),
                applicable_roles=q.get("applicable_roles", ["all"])
            )
        )

    assign_id = existing_assignment.assignment_id if existing_assignment else f"mospi-baseline-{div_code.lower()}-{role_cat}-v1"

    return BaselineAssessmentOut(
        assessment_id=assign_id,
        title=f"Baseline Diagnostic Assessment ({div_title} — {desig_title})",
        instructions=f"Complete this calibrated diagnostic evaluation tailored specifically for your role as {desig_title} in {div_title}. Contains {bp_common_count} Common Core statistical questions and {bp_role_count} Designation/Role-Specific questions.",
        total_questions=len(questions_out),
        time_limit_mins=20,
        questions=questions_out,
        officer_designation=desig_title,
        role_tier=role_cat,
        common_core_count=bp_common_count,
        role_specific_count=bp_role_count,
        assessed_competencies=assessed_comps
    )

def evaluate_baseline_submission(user_id: int, submission: BaselineAssessmentSubmit, db: Session) -> BaselineAssessmentResultOut:
    user = db.query(User).filter(User.id == user_id).first()
    role_meta = resolve_role_benchmarks(user.department if user else "", user.designation if user else "")
    div_code = role_meta.get("division_code", "GENERAL")

    answers_map = {a.question_id: a.selected_option.strip().upper() for a in submission.answers}

    from ..models.models import BaselineAssignment
    assignment = db.query(BaselineAssignment).filter(
        BaselineAssignment.user_id == user_id
    ).order_by(BaselineAssignment.assigned_at.desc()).first()

    if assignment and assignment.status == "submitted":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Duplicate submission rejected: This baseline assessment assignment has already been submitted."
        )

    all_known_questions = {}
    for q in COMMON_CORE_QUESTIONS:
        all_known_questions[q["id"]] = dict(q)
    for q_list in ROLE_SPECIFIC_QUESTION_BANK.values():
        for q in q_list:
            all_known_questions[q["id"]] = dict(q)
    for q_list in DEPARTMENT_BASELINE_BANK.values():
        for q in q_list:
            all_known_questions[q["id"]] = dict(q)

    target_questions = []
    if assignment:
        try:
            assigned_ids = json.loads(assignment.assigned_question_ids)
            assigned_id_set = set(assigned_ids)

            submitted_ids = set(answers_map.keys())
            if not submitted_ids.issubset(assigned_id_set):
                unassigned = submitted_ids - assigned_id_set
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid submission: Question IDs {list(unassigned)} were not assigned to your baseline assessment."
                )

            for qid in assigned_ids:
                if qid in all_known_questions:
                    target_questions.append(all_known_questions[qid])
        except HTTPException:
            raise
        except Exception as e:
            logger.warning(f"[BaselineSubmission] Error parsing assignment assigned_question_ids: {e}")

    if not target_questions:
        target_questions = USER_BASELINE_CACHE.get(user_id)

    if not target_questions:
        for a in submission.answers:
            if a.question_id in all_known_questions:
                target_questions.append(all_known_questions[a.question_id])

        if not target_questions:
            target_questions = DEPARTMENT_BASELINE_BANK.get(div_code, BASELINE_QUESTIONS)

    total_questions = len(target_questions)
    total_correct = 0
    domain_correct = {}
    domain_totals = {}
    comp_totals = {}
    comp_correct = {}

    for q in target_questions:
        qid = q["id"]
        domain = q["domain"]
        comp_code = q["competency_code"]

        domain_totals[domain] = domain_totals.get(domain, 0) + 1
        comp_totals[comp_code] = comp_totals.get(comp_code, 0) + 1

        user_ans = answers_map.get(qid, "")
        correct_ans = q.get("correct_option", "A")

        is_correct = (user_ans == correct_ans)
        if is_correct:
            total_correct += 1
            domain_correct[domain] = domain_correct.get(domain, 0) + 1
            comp_correct[comp_code] = comp_correct.get(comp_code, 0) + 1

    competency_scores = {}
    for code, total_q in comp_totals.items():
        corr = comp_correct.get(code, 0)
        score_pct = (corr / total_q) * 100.0
        if score_pct == 100.0:
            competency_scores[code] = 75.0
        elif score_pct > 0.0:
            competency_scores[code] = 60.0
        else:
            competency_scores[code] = 40.0

    overall_score = round((total_correct / total_questions) * 100.0, 1) if total_questions > 0 else 0.0

    domain_scores = {}
    for d, tot in domain_totals.items():
        corr = domain_correct.get(d, 0)
        domain_scores[d] = round((corr / tot) * 100.0, 1)

    all_competencies = db.query(Competency).all()
    comp_code_map = {c.code: c for c in all_competencies}

    for code, score in competency_scores.items():
        comp_obj = comp_code_map.get(code)
        if comp_obj:
            uc = db.query(UserCompetency).filter(
                UserCompetency.user_id == user_id,
                UserCompetency.competency_id == comp_obj.id
            ).first()
            if not uc:
                uc = UserCompetency(
                    user_id=user_id,
                    competency_id=comp_obj.id,
                    current_level=score,
                    assessment_source="baseline_assessment",
                    last_assessed_at=datetime.utcnow()
                )
                db.add(uc)
            else:
                uc.current_level = score
                uc.assessment_source = "baseline_assessment"
                uc.last_assessed_at = datetime.utcnow()

            hist = LearningProgressHistory(
                user_id=user_id,
                competency_id=comp_obj.id,
                event_type="baseline_assessment",
                previous_score=0.0,
                new_score=score,
                delta=score,
                created_at=datetime.utcnow()
            )
            db.add(hist)

    if assignment:
        assignment.status = "submitted"
        assignment.submitted_at = datetime.utcnow()
        assignment.score = overall_score

    db.commit()

    div_title = user.department if user and user.department else "MoSPI"
    desig_title = user.designation if user and user.designation else "Statistical Officer"

    feedback = (
        f"Designation-based baseline assessment for {desig_title} ({div_title}) completed! "
        f"You achieved an overall score of {overall_score}% ({total_correct}/{total_questions} correct). "
        f"Your scores across Common Core and Role-Specific questions have initialized {len(competency_scores)} statistical competencies against official {div_title} benchmarks. "
        f"Explore your prioritized competency gap analysis to begin targeted learning."
    )

    return BaselineAssessmentResultOut(
        overall_score=overall_score,
        total_correct=total_correct,
        total_questions=total_questions,
        domain_scores=domain_scores,
        competency_scores=competency_scores,
        initialized_competencies_count=len(competency_scores),
        feedback_summary=feedback
    )

def validate_mcq_quality(q: Dict[str, Any], seen_texts: Set[str]) -> bool:
    q_text = (q.get("question_text") or "").strip()
    expl = (q.get("explanation") or "").strip()
    if not q_text or not expl:
        return False
    opts = [q.get("option_a"), q.get("option_b"), q.get("option_c"), q.get("option_d")]
    if any(not opt or not str(opt).strip() for opt in opts):
        return False
    clean_opts = [str(opt).strip().lower() for opt in opts]
    if len(set(clean_opts)) < 4:
        return False
    corr = (q.get("correct_option") or "").strip().upper()
    if corr not in ["A", "B", "C", "D"]:
        return False
    q_norm = q_text.lower()
    if q_norm in seen_texts:
        return False
    seen_texts.add(q_norm)
    return True

async def create_ai_quiz(request: QuizGenerateRequest, user_id: int, db: Session) -> QuizOut:
    source_text = ""
    doc_obj = None
    target_comp = None

    if request.resource_id:
        res_obj = db.query(LearningResource).filter(LearningResource.id == request.resource_id).first()
        if res_obj:
            source_text = f"Title: {res_obj.title}\nSource: {res_obj.source}\nType: {res_obj.resource_type}\nDescription: {res_obj.description}"
            if res_obj.competency_mappings:
                target_comp = res_obj.competency_mappings[0].competency
    elif request.document_id:
        doc_obj = db.query(Document).filter(
            Document.id == request.document_id,
            Document.user_id == user_id
        ).first()
        if not doc_obj:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Uploaded document not found.")
        source_text = doc_obj.extracted_text
    elif request.custom_text:
        source_text = request.custom_text
    else:
        source_text = f"Official statistical concepts and guidelines on {request.topic}"

    if request.competency_id:
        target_comp = db.query(Competency).filter(Competency.id == request.competency_id).first()
    if not target_comp:
        topic_lower = request.topic.lower()
        if "python" in topic_lower or "comput" in topic_lower or "data" in topic_lower:
            target_comp = db.query(Competency).filter(Competency.code == "STAT_COMPUTE").first()
        elif "account" in topic_lower or "gdp" in topic_lower or "gva" in topic_lower:
            target_comp = db.query(Competency).filter(Competency.code == "STAT_NAT_ACC").first()
        elif "price" in topic_lower or "cpi" in topic_lower or "iip" in topic_lower:
            target_comp = db.query(Competency).filter(Competency.code == "STAT_PRICE_IND").first()
        else:
            target_comp = db.query(Competency).filter(Competency.code == "STAT_SURVEY").first()

    raw_questions = await generate_mcqs_from_document_async(
        text=source_text,
        topic=request.topic,
        num_questions=request.num_questions,
        difficulty=request.difficulty,
        competency_code=target_comp.code if target_comp else None
    )

    # MCQ Quality Validation
    seen_q_texts: Set[str] = set()
    validated_questions = []
    for q_data in raw_questions:
        if validate_mcq_quality(q_data, seen_q_texts):
            validated_questions.append(q_data)

    if len(validated_questions) < request.num_questions:
        validated_questions = raw_questions[:request.num_questions]

    # Detect generation method
    gen_method = "DETERMINISTIC_FALLBACK"
    purpose_val = getattr(request, "purpose", "SELF_ASSESSMENT") or "SELF_ASSESSMENT"

    blueprint_dict = {
        "total_questions": len(validated_questions),
        "mcq_count": len(validated_questions),
        "difficulty_distribution": {request.difficulty: len(validated_questions)},
        "competency_coverage": [target_comp.code if target_comp else "STAT_SURVEY"],
        "source_document_references": [request.document_id] if request.document_id else [],
        "requested_difficulty": request.difficulty,
        "purpose": purpose_val
    }

    new_quiz = Quiz(
        user_id=user_id,
        document_id=request.document_id,
        competency_id=target_comp.id if target_comp else None,
        title=f"AI Quiz: {request.topic}",
        topic=request.topic,
        difficulty=request.difficulty,
        total_questions=len(validated_questions),
        time_limit_mins=max(5, len(validated_questions) * 3),
        purpose=purpose_val,
        blueprint_metadata=json.dumps(blueprint_dict),
        generation_method=gen_method
    )
    db.add(new_quiz)
    db.commit()
    db.refresh(new_quiz)

    q_out_list: List[QuizQuestionOut] = []
    for idx, q_data in enumerate(validated_questions):
        source_ref = f"doc:{request.document_id}#chunk:{idx}" if request.document_id else f"topic:{request.topic}"
        q_record = QuizQuestion(
            quiz_id=new_quiz.id,
            question_text=q_data["question_text"],
            option_a=q_data["option_a"],
            option_b=q_data["option_b"],
            option_c=q_data["option_c"],
            option_d=q_data["option_d"],
            correct_option=q_data["correct_option"],
            explanation=q_data["explanation"],
            competency_code=target_comp.code if target_comp else "STAT_SURVEY",
            difficulty=request.difficulty,
            source_reference=source_ref,
            generation_method=gen_method,
            competency_mapping_method=doc_obj.mapping_method if doc_obj and hasattr(doc_obj, 'mapping_method') and doc_obj.mapping_method else "PLATFORM_HEURISTIC"
        )
        db.add(q_record)
        db.commit()
        db.refresh(q_record)

        q_out_list.append(
            QuizQuestionOut(
                id=q_record.id,
                question_text=q_record.question_text,
                option_a=q_record.option_a,
                option_b=q_record.option_b,
                option_c=q_record.option_c,
                option_d=q_record.option_d,
                difficulty=q_record.difficulty,
                source_reference=q_record.source_reference,
                generation_method=q_record.generation_method,
                competency_mapping_method=q_record.competency_mapping_method
            )
        )

    return QuizOut(
        id=new_quiz.id,
        title=new_quiz.title,
        topic=new_quiz.topic,
        difficulty=new_quiz.difficulty,
        total_questions=new_quiz.total_questions,
        time_limit_mins=new_quiz.time_limit_mins,
        purpose=new_quiz.purpose,
        generation_method=new_quiz.generation_method,
        blueprint_metadata=new_quiz.blueprint_metadata,
        created_at=new_quiz.created_at,
        questions=q_out_list
    )

async def evaluate_quiz_submission(
    quiz_id: int,
    user_id: int,
    submission: QuizSubmitRequest,
    db: Session
) -> QuizAttemptResultOut:
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Quiz not found.")

    existing_attempt = db.query(QuizAttempt).filter(
        QuizAttempt.quiz_id == quiz_id,
        QuizAttempt.user_id == user_id,
        QuizAttempt.status.in_(["SUBMITTED", "EVALUATED"])
    ).first()

    if existing_attempt:
        # Re-submission protection
        questions = db.query(QuizQuestion).filter(QuizQuestion.quiz_id == quiz_id).all()
        q_results = [
            QuestionResultDetail(
                question_id=q.id,
                question_text=q.question_text,
                user_selected=q.correct_option,
                correct_option=q.correct_option,
                is_correct=True,
                explanation=q.explanation,
                source_reference=q.source_reference
            ) for q in questions
        ]
        comp_obj = db.query(Competency).filter(Competency.id == quiz.competency_id).first() if quiz.competency_id else None
        return QuizAttemptResultOut(
            attempt_id=existing_attempt.id,
            quiz_id=quiz.id,
            quiz_title=quiz.title,
            score=existing_attempt.score,
            total_correct=existing_attempt.total_correct,
            total_questions=existing_attempt.total_questions,
            status="EVALUATED",
            feedback_method=existing_attempt.feedback_method or "Deterministic Pedagogical Feedback",
            competency_name=comp_obj.name if comp_obj else quiz.topic,
            competency_score_before=existing_attempt.competency_score_before,
            competency_score_after=existing_attempt.competency_score_after,
            competency_delta=0.0,
            ai_qualitative_feedback=existing_attempt.ai_qualitative_feedback or "Attempt already evaluated.",
            question_results=q_results,
            completed_at=existing_attempt.completed_at
        )

    questions = db.query(QuizQuestion).filter(QuizQuestion.quiz_id == quiz_id).all()
    if not questions:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No questions in this quiz.")

    answers_map = {a.question_id: a.selected_option.strip().upper() for a in submission.answers}
    total_correct = 0
    question_results: List[QuestionResultDetail] = []

    for q in questions:
        user_sel = answers_map.get(q.id, "")
        is_corr = (user_sel == q.correct_option)
        if is_corr:
            total_correct += 1

        question_results.append(
            QuestionResultDetail(
                question_id=q.id,
                question_text=q.question_text,
                user_selected=user_sel if user_sel else "None",
                correct_option=q.correct_option,
                is_correct=is_corr,
                explanation=q.explanation,
                source_reference=q.source_reference
            )
        )

    score_pct = round((total_correct / len(questions)) * 100.0, 1)

    comp_obj = db.query(Competency).filter(Competency.id == quiz.competency_id).first() if quiz.competency_id else None
    before_score = 42.0
    after_score = 42.0
    delta = 0.0

    attempt_rec = QuizAttempt(
        quiz_id=quiz_id,
        user_id=user_id,
        score=score_pct,
        total_correct=total_correct,
        total_questions=len(questions),
        competency_id=comp_obj.id if comp_obj else None,
        competency_score_before=0.0,
        competency_score_after=0.0,
        competency_delta=0.0,
        status="EVALUATED",
        feedback_method="Deterministic Pedagogical Feedback",
        evidence_key=f"quiz-attempt-{quiz_id}-user-{user_id}",
        completed_at=datetime.utcnow()
    )
    db.add(attempt_rec)
    db.flush()

    if comp_obj:
        user_comp = db.query(UserCompetency).filter(
            UserCompetency.user_id == user_id,
            UserCompetency.competency_id == comp_obj.id
        ).first()
        before_score = user_comp.current_level if user_comp else 0.0

        from .learning_adaptive_service import process_learning_evidence
        ev_res = process_learning_evidence(
            db=db,
            user_id=user_id,
            competency_id=comp_obj.id,
            evidence_type="QUIZ_ATTEMPT",
            evidence_key=f"quiz-attempt-{attempt_rec.id}",
            score_percentage=score_pct,
            difficulty=quiz.difficulty or "Intermediate",
            explanation=f"Completed quiz: {quiz.title}"
        )

        after_score = ev_res["new_score"]
        delta = ev_res["delta"]

        attempt_rec.competency_score_before = before_score
        attempt_rec.competency_score_after = after_score
        attempt_rec.competency_delta = delta
        db.commit()

    mistakes = [
        {
            "question_text": q.question_text,
            "user_selected": q.user_selected,
            "correct_option": q.correct_option,
            "explanation": q.explanation
        }
        for q in question_results if not q.is_correct
    ]

    feedback = await generate_grok_quiz_feedback(
        quiz_title=quiz.title,
        topic=quiz.topic,
        score_pct=score_pct,
        total_correct=total_correct,
        total_questions=len(questions),
        competency_name=comp_obj.name if comp_obj else quiz.topic,
        before_score=before_score,
        after_score=after_score,
        delta=delta,
        mistakes=mistakes
    )

    attempt_rec.ai_qualitative_feedback = feedback
    attempt_rec.feedback_method = "Deterministic Pedagogical Feedback"
    db.commit()
    db.refresh(attempt_rec)

    return QuizAttemptResultOut(
        attempt_id=attempt_rec.id,
        quiz_id=quiz.id,
        quiz_title=quiz.title,
        score=score_pct,
        total_correct=total_correct,
        total_questions=len(questions),
        status=attempt_rec.status,
        feedback_method=attempt_rec.feedback_method,
        competency_name=comp_obj.name if comp_obj else quiz.topic,
        competency_score_before=before_score,
        competency_score_after=after_score,
        competency_delta=delta,
        ai_qualitative_feedback=feedback,
        question_results=question_results,
        completed_at=attempt_rec.completed_at
    )

def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"PDF extraction error: {e}")

def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"DOCX extraction error: {e}")

def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"PPTX extraction error: {e}")

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="replace").strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Text read error: {e}")

def process_uploaded_document(filename: str, file_bytes: bytes) -> str:
    ext = filename.split(".")[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext in ["docx", "doc"]:
        return extract_text_from_docx(file_bytes)
    elif ext in ["pptx", "ppt"]:
        return extract_text_from_pptx(file_bytes)
    else:
        return extract_text_from_txt(file_bytes)

def get_user_progress_summary(user_id: int, db: Session) -> ProgressSummaryOut:
    user = db.query(User).filter(User.id == user_id).first()
    all_competencies = db.query(Competency).all()
    user_comps = {uc.competency_id: uc for uc in db.query(UserCompetency).filter(UserCompetency.user_id == user_id).all()}
    history_records = db.query(LearningProgressHistory).filter(
        LearningProgressHistory.user_id == user_id
    ).order_by(LearningProgressHistory.created_at.desc()).all()
    attempts = db.query(QuizAttempt).filter(QuizAttempt.user_id == user_id).all()

    cards: List[CompetencyProgressCard] = []
    total_gain_sum = 0.0
    total_score = 0.0

    for comp in all_competencies:
        uc = user_comps.get(comp.id)
        current = uc.current_level if uc else 0.0
        total_score += current

        comp_hist = [h for h in history_records if h.competency_id == comp.id]
        if comp_hist:
            initial = comp_hist[-1].previous_score
            gain = max(0.0, round(current - initial, 1))
        else:
            initial = current
            gain = 0.0

        total_gain_sum += gain

        if current >= comp.required_level:
            status_str = "Mastered"
        elif gain > 0:
            status_str = "Improving"
        else:
            status_str = "Needs Attention"

        cards.append(
            CompetencyProgressCard(
                competency_id=comp.id,
                code=comp.code,
                name=comp.name,
                domain=comp.domain,
                initial_score=initial,
                current_score=current,
                required_benchmark=comp.required_level,
                total_gain=gain,
                status=status_str
            )
        )

    overall_readiness = round(total_score / len(all_competencies), 1) if all_competencies else 0.0
    avg_quiz_score = round(sum(a.score for a in attempts) / len(attempts), 1) if attempts else 0.0

    recent_events_out: List[ProgressEventOut] = []
    for h in history_records[:10]:
        comp_obj = db.query(Competency).filter(Competency.id == h.competency_id).first()
        comp_name = comp_obj.name if comp_obj else "General Statistics"
        comp_domain = comp_obj.domain if comp_obj else "Official Statistics"

        recent_events_out.append(
            ProgressEventOut(
                id=h.id,
                competency_id=h.competency_id,
                competency_name=comp_name,
                domain=comp_domain,
                event_type=h.event_type,
                previous_score=h.previous_score,
                new_score=h.new_score,
                delta=h.delta,
                created_at=h.created_at
            )
        )

    return ProgressSummaryOut(
        user_id=user.id if user else 0,
        user_name=user.full_name if user else "Officer",
        designation=user.designation if user else "Statistical Officer",
        department=user.department if user else "MoSPI",
        overall_readiness_score=overall_readiness,
        total_learning_gain=round(total_gain_sum / len(all_competencies), 1) if all_competencies else 0.0,
        quizzes_completed=len(attempts),
        average_quiz_score=avg_quiz_score,
        competency_breakdown=cards,
        recent_events=recent_events_out,
        progress_events=recent_events_out
    )


def get_final_interview_readiness(user_id: int, db: Session) -> FinalInterviewReadiness:
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        return FinalInterviewReadiness(
            eligible=False, readiness_score=0.0, competencies_to_assess=[], message="User not found."
        )

    competencies = db.query(Competency).all()
    user_competencies = {
        uc.competency_id: uc
        for uc in db.query(UserCompetency).filter(UserCompetency.user_id == user_id).all()
    }

    competency_cards = []
    total_score = 0.0

    for competency in competencies:
        user_comp = user_competencies.get(competency.id)
        current_score = user_comp.current_level if user_comp else 0.0
        gap = max(0.0, competency.required_level - current_score)
        total_score += current_score

        competency_cards.append(
            FinalInterviewCompetency(
                competency_id=competency.id,
                code=competency.code,
                name=competency.name,
                domain=competency.domain,
                current_score=round(current_score, 1),
                required_benchmark=round(competency.required_level, 1),
                gap=round(gap, 1)
            )
        )

    readiness_score = round(total_score / len(competencies), 1) if competencies else 0.0
    competency_cards.sort(key=lambda item: item.gap, reverse=True)

    return FinalInterviewReadiness(
        eligible=True,
        readiness_score=readiness_score,
        competencies_to_assess=competency_cards,
        message="You are ready for your final AI interview."
    )

async def generate_interview_questions(user_id: int, db: Session, num_questions: int = 5):
    readiness = get_final_interview_readiness(user_id, db)
    if not readiness.eligible:
        return {"eligible": False, "questions": [], "message": readiness.message}

    competencies = [competency.model_dump() if hasattr(competency, 'model_dump') else competency.dict() for competency in readiness.competencies_to_assess]
    questions = await generate_final_interview_questions(competencies=competencies, num_questions=num_questions)

    return {
        "eligible": True,
        "readiness_score": readiness.readiness_score,
        "questions": questions,
        "message": "Final AI interview questions generated successfully."
    }
