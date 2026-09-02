import io
import re
import hashlib
import logging
from typing import List, Dict, Any, Optional
from datetime import datetime
from fastapi import HTTPException, status
from sqlalchemy.orm import Session
from pypdf import PdfReader
import docx
from pptx import Presentation

from ..models.models import Competency, Document, ContentChunk

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {"pdf", "docx", "doc", "pptx", "ppt", "txt", "md"}
MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024  # 25 MB
MIN_TEXT_THRESHOLD_CHARS = 30

COMPETENCY_KEYWORD_RULES = {
    "STAT_SURVEY": ["survey", "sampling", "strata", "multistage", "nsso", "sample size", "enumeration", "plfs"],
    "STAT_NAT_ACC": ["national accounts", "gdp", "gva", "sna 2008", "gross domestic product", "gross value added", "deflator", "macroeconomic"],
    "STAT_COMPUTE": ["python", "pandas", "numpy", "r", "stata", "sql", "microdata", "data science", "computing", "pipeline"],
    "STAT_PRICE_IND": ["cpi", "wpi", "iip", "price index", "laspeyres", "paasche", "inflation", "index number"],
    "STAT_LABOUR": ["labour", "plfs", "employment", "unemployment", "lfpr", "wpr", "upss", "cws"],
    "STAT_DATA_GOV": ["esankhyiki", "metadata", "governance", "fair", "data architecture", "anonymization"],
    "STAT_QUALITY": ["quality", "audit", "nqaf", "fundamental principles", "imputation", "supervision"],
    "STAT_VIZ_COMM": ["visualization", "sdg", "dashboard", "storytelling", "dissemination", "press release"],
    "STAT_IND_AGRI": ["asi", "annual survey of industries", "enterprise", "factory frame", "npc", "nic"]
}

def normalize_text(text: str) -> str:
    if not text:
        return ""
    # Normalize Windows/Mac line endings
    clean = text.replace("\r\n", "\n").replace("\r", "\n")
    # Replace multiple consecutive spaces within lines, preserving paragraph breaks
    paragraphs = clean.split("\n\n")
    normalized_paras = []
    for para in paragraphs:
        p_clean = re.sub(r"[ \t]+", " ", para).strip()
        if p_clean:
            normalized_paras.append(p_clean)
    return "\n\n".join(normalized_paras)

def extract_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="replace").strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Text extraction failed: {e}")

def extract_pdf(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"PDF extraction failed: {e}")

def extract_docx(file_bytes: bytes) -> str:
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"DOCX extraction failed: {e}")

def extract_pptx(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"PPTX extraction failed: {e}")

def process_and_validate_file(filename: str, file_bytes: bytes) -> Dict[str, Any]:
    if not filename or "." not in filename:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="File has no extension.")

    ext = filename.split(".")[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file format: '.{ext}'. Supported formats are: {', '.join(sorted(ALLOWED_EXTENSIONS))}."
        )

    if len(file_bytes) == 0:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Uploaded file is empty.")

    if len(file_bytes) > MAX_FILE_SIZE_BYTES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"File size exceeds maximum allowed limit of {MAX_FILE_SIZE_BYTES // (1024 * 1024)}MB."
        )

    content_hash = hashlib.sha256(file_bytes).hexdigest()

    if ext == "pdf":
        raw_text = extract_pdf(file_bytes)
    elif ext in ["docx", "doc"]:
        raw_text = extract_docx(file_bytes)
    elif ext in ["pptx", "ppt"]:
        raw_text = extract_pptx(file_bytes)
    else:
        raw_text = extract_txt(file_bytes)

    norm_text = normalize_text(raw_text)

    if len(norm_text) < MIN_TEXT_THRESHOLD_CHARS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Extracted content is too short ({len(norm_text)} chars). Minimum threshold is {MIN_TEXT_THRESHOLD_CHARS} chars."
        )

    return {
        "ext": ext,
        "content_hash": content_hash,
        "extracted_text": norm_text,
        "character_count": len(norm_text),
        "extraction_status": "SUCCESS"
    }

def chunk_text(
    text: str,
    document_id: int,
    max_chunk_size: int = 1000,
    min_chunk_size: int = 100,
    overlap: int = 100
) -> List[Dict[str, Any]]:
    if not text:
        return []

    paragraphs = text.split("\n\n")
    chunks: List[Dict[str, Any]] = []
    curr_text = ""
    chunk_idx = 0

    for para in paragraphs:
        if len(curr_text) + len(para) + 2 <= max_chunk_size:
            curr_text = f"{curr_text}\n\n{para}".strip()
        else:
            if len(curr_text) >= min_chunk_size:
                chunks.append({
                    "chunk_index": chunk_idx,
                    "chunk_text": curr_text,
                    "character_count": len(curr_text),
                    "token_count_approx": max(1, len(curr_text) // 5),
                    "document_id": document_id
                })
                chunk_idx += 1
                # Preserve overlap if available
                overlap_text = curr_text[-overlap:] if len(curr_text) >= overlap else curr_text
                curr_text = f"{overlap_text}\n\n{para}".strip()
            else:
                curr_text = f"{curr_text}\n\n{para}".strip()

    if curr_text.strip():
        chunks.append({
            "chunk_index": chunk_idx,
            "chunk_text": curr_text.strip(),
            "character_count": len(curr_text.strip()),
            "token_count_approx": max(1, len(curr_text.strip()) // 5),
            "document_id": document_id
        })

    return chunks

def map_content_to_competencies(text: str, db: Session) -> Dict[str, Any]:
    text_lower = text.lower()
    scores: Dict[str, int] = {}

    for code, keywords in COMPETENCY_KEYWORD_RULES.items():
        score = sum(text_lower.count(kw) for kw in keywords)
        if score > 0:
            scores[code] = score

    if not scores:
        default_comp = db.query(Competency).filter(Competency.code == "STAT_SURVEY").first()
        return {
            "competency_id": default_comp.id if default_comp else None,
            "competency_code": default_comp.code if default_comp else "STAT_SURVEY",
            "competency_name": default_comp.name if default_comp else "Survey Methodology & Sampling Design",
            "mapping_confidence": 0.5,
            "mapping_method": "PLATFORM_HEURISTIC"
        }

    top_code = max(scores.items(), key=lambda x: x[1])[0]
    matched_comp = db.query(Competency).filter(Competency.code == top_code).first()

    total_matches = sum(scores.values())
    confidence = round(min(0.95, 0.5 + (scores[top_code] / (total_matches + 2))), 2)

    return {
        "competency_id": matched_comp.id if matched_comp else None,
        "competency_code": matched_comp.code if matched_comp else top_code,
        "competency_name": matched_comp.name if matched_comp else top_code,
        "mapping_confidence": confidence,
        "mapping_method": "PLATFORM_HEURISTIC"
    }
